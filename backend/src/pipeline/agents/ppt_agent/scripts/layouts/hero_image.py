"""hero_image layout — text on left 55%, generated image on right 45%.

Used for: high-impact main slides (solution overview, key differentiator).
Only assigned when SlideBlueprint.image_source == 'generated' and a
generated_image_path has been attached by the image_generator node.

Reads from content:
  - title + key_bullets (or simplified_bullets if density_verdict='simplify').

Reads from blueprint:
  - generated_image_path — the FLUX-produced PNG to embed on the right panel.

Falls back to bullets_only if generated_image_path is missing (should
not happen in production — indicates a planner/image-gen bug).
"""

from __future__ import annotations

from pptx.util import Inches

from pipeline.agents.ppt_agent.schema import (
    ContentSource,
    DensityVerdict,
    SalesPPTOutput,
    SlideBlueprint,
)
from pipeline.agents.ppt_agent.scripts.render.image_helpers import fit_image_dims
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    BODY_TOP,
    MARGIN_BOTTOM,
    MARGIN_LEFT,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    add_bullets,
    add_footer,
    add_logo,
    add_title,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design

# Split: text 55% left, image 45% right.
_TEXT_WIDTH_RATIO = 0.55
_IMAGE_WIDTH_RATIO = 0.45
_GAP = Inches(0.3)


def _get_content(
    content: SalesPPTOutput, blueprint: SlideBlueprint,
) -> tuple[str, list[str]]:
    """Extract title and bullets based on content_source."""
    src = blueprint.content_source

    if src == ContentSource.MAIN_SLIDE:
        slide = content.slides[blueprint.content_index]
        bullets = slide.key_bullets
        if (
            blueprint.density_verdict == DensityVerdict.SIMPLIFY
            and blueprint.simplified_bullets
        ):
            bullets = blueprint.simplified_bullets
        return slide.title, bullets

    if src == ContentSource.DECK_WHY_AZIRO:
        return "Why Aziro", content.why_aziro

    if src == ContentSource.DECK_DIFFERENTIATORS:
        return "Differentiators", content.differentiators

    # Fallback.
    return src.value.replace("_", " ").title(), []


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render a hero_image slide — text left, image right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    title, bullets = _get_content(content, blueprint)

    usable_width = SLIDE_WIDTH - MARGIN_LEFT * 2
    text_width = int(usable_width * _TEXT_WIDTH_RATIO) - _GAP

    # Title and bullets on the left panel.
    add_title(slide, title, design, width=text_width)
    add_bullets(slide, bullets, design, width=text_width)

    # Image on the right panel.
    if blueprint.generated_image_path:
        img_max_w = int(usable_width * _IMAGE_WIDTH_RATIO)
        img_max_h = SLIDE_HEIGHT - BODY_TOP - MARGIN_BOTTOM
        img_left = MARGIN_LEFT + int(usable_width * _TEXT_WIDTH_RATIO)

        img_w, img_h = fit_image_dims(
            blueprint.generated_image_path, img_max_w, img_max_h,
        )
        # Centre the image vertically in the right panel.
        img_top = BODY_TOP + (img_max_h - img_h) // 2

        slide.shapes.add_picture(
            blueprint.generated_image_path,
            left=img_left,
            top=img_top,
            width=img_w,
            height=img_h,
        )

    add_logo(slide, design)
    add_footer(slide, design)
