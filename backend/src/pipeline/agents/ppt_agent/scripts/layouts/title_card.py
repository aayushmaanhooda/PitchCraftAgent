"""title_card layout — cover / hero with optional right-side image.

Used for:
  - The auto-emitted cover slide (content_source == deck_title). When a
    generated_image_path is present the title occupies the left 55% and
    the image fills the right 45%. Otherwise the title is centred.
  - Any main_slide the planner flags as title_card (subtitle = purpose).

Styling: heading_font + primary_color for the hero title, body_font +
secondary_color for the subtitle. Logo + footer from the Design.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from pipeline.agents.ppt_agent.schema import ContentSource, SalesPPTOutput, SlideBlueprint
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    CONTENT_WIDTH,
    MARGIN_LEFT,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    add_footer,
    add_logo,
    add_textbox,
    fit_heading_box,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.render.image_helpers import fit_image_dims
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design


def _get_content(
    content: SalesPPTOutput, blueprint: SlideBlueprint,
) -> tuple[str, str | None]:
    """Extract title and optional subtitle based on content_source."""
    src = blueprint.content_source

    if src == ContentSource.DECK_TITLE:
        return content.deck_title, content.proposed_solution_theme

    if src == ContentSource.MAIN_SLIDE:
        slide = content.slides[blueprint.content_index]
        return slide.title, slide.purpose

    return src.value.replace("_", " ").title(), None


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render a title_card slide — text left, image right when available."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    title, subtitle = _get_content(content, blueprint)
    has_image = bool(blueprint.generated_image_path)

    if has_image:
        # Text left 55%, image right 45%.
        text_left = MARGIN_LEFT
        text_width_in = 6.8
        text_width = Inches(text_width_in)

        title_font_pt, title_height_in = fit_heading_box(
            title,
            width_in=text_width_in,
            base_font_pt=36,
            min_font_pt=22,
            max_lines=4,
        )

        # Vertically centre the title + subtitle block on the slide.
        subtitle_gap_in = 0.25
        subtitle_height_in = 1.3 if subtitle else 0.0
        block_height_in = title_height_in + (
            subtitle_gap_in + subtitle_height_in if subtitle else 0
        )
        text_top_in = max(0.8, (7.5 - block_height_in) / 2)
        text_top = Inches(text_top_in)

        add_textbox(
            slide,
            left=text_left,
            top=text_top,
            width=text_width,
            height=Inches(title_height_in),
            text=title,
            font_name=design.heading_font,
            font_size_pt=title_font_pt,
            color_hex=design.primary_color,
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        if subtitle:
            add_textbox(
                slide,
                left=text_left,
                top=Inches(text_top_in + title_height_in + subtitle_gap_in),
                width=text_width,
                height=Inches(subtitle_height_in),
                text=subtitle,
                font_name=design.body_font,
                font_size_pt=18,
                color_hex=design.secondary_color,
                alignment=PP_ALIGN.LEFT,
            )

        # Image on the right.
        img_region_left = Inches(7.6)
        img_region_top = Inches(1.2)
        img_max_w = Inches(5.2)
        img_max_h = Inches(5.0)
        try:
            w, h = fit_image_dims(
                blueprint.generated_image_path, img_max_w, img_max_h,
            )
            img_left = img_region_left + (img_max_w - w) // 2
            img_top = img_region_top + (img_max_h - h) // 2
            slide.shapes.add_picture(
                blueprint.generated_image_path,
                left=img_left, top=img_top, width=w, height=h,
            )
        except Exception:
            pass

    else:
        # Centred hero title on an empty canvas.
        content_width_in = 13.333 - 0.6 - 0.6
        title_font_pt, title_height_in = fit_heading_box(
            title,
            width_in=content_width_in,
            base_font_pt=40,
            min_font_pt=26,
            max_lines=3,
        )
        subtitle_gap_in = 0.3
        subtitle_height_in = 1.0 if subtitle else 0.0
        block_height_in = title_height_in + (
            subtitle_gap_in + subtitle_height_in if subtitle else 0
        )
        title_top_in = max(1.5, (7.5 - block_height_in) / 2)

        add_textbox(
            slide,
            left=MARGIN_LEFT,
            top=Inches(title_top_in),
            width=CONTENT_WIDTH,
            height=Inches(title_height_in),
            text=title,
            font_name=design.heading_font,
            font_size_pt=title_font_pt,
            color_hex=design.primary_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            add_textbox(
                slide,
                left=MARGIN_LEFT,
                top=Inches(title_top_in + title_height_in + subtitle_gap_in),
                width=CONTENT_WIDTH,
                height=Inches(subtitle_height_in),
                text=subtitle,
                font_name=design.body_font,
                font_size_pt=20,
                color_hex=design.secondary_color,
                alignment=PP_ALIGN.CENTER,
            )

    add_logo(slide, design)
    add_footer(slide, design)
