"""summary_solution layout — two stacked sections + right-side image.

Used ONLY for ContentSource.DECK_SUMMARY_SOLUTION.

Layout:
  Title:  "Overview"
  Left 55%:
    Section "Summary"  — executive_summary prose
    Section "Solution" — proposed_solution_theme line
  Right 45%: deterministic FLUX image if generated_image_path is set.
"""

from __future__ import annotations

from pptx.util import Inches

from pipeline.agents.ppt_agent.schema import SalesPPTOutput, SlideBlueprint
from pipeline.agents.ppt_agent.scripts.render.image_helpers import fit_image_dims
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    BODY_LEFT,
    BODY_TOP,
    add_footer,
    add_logo,
    add_textbox,
    add_title,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design

_TEXT_WIDTH = Inches(6.8)
_SUBHEAD_HEIGHT = Inches(0.45)
_SECTION_GAP = Inches(0.3)


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render a merged Summary + Solution slide with a right-side image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    add_title(slide, "Overview", design)

    # ── Section 1: Summary ──────────────────────────────────────────
    section_top = BODY_TOP
    add_textbox(
        slide,
        left=BODY_LEFT,
        top=section_top,
        width=_TEXT_WIDTH,
        height=_SUBHEAD_HEIGHT,
        text="Summary",
        font_name=design.heading_font,
        font_size_pt=20,
        color_hex=design.accent_color,
        bold=True,
    )
    summary_top = section_top + _SUBHEAD_HEIGHT
    summary_height = Inches(2.2)
    add_textbox(
        slide,
        left=BODY_LEFT,
        top=summary_top,
        width=_TEXT_WIDTH,
        height=summary_height,
        text=content.executive_summary,
        font_name=design.body_font,
        font_size_pt=14,
        color_hex=design.text_color,
    )

    # ── Section 2: Solution ─────────────────────────────────────────
    section2_top = summary_top + summary_height + _SECTION_GAP
    add_textbox(
        slide,
        left=BODY_LEFT,
        top=section2_top,
        width=_TEXT_WIDTH,
        height=_SUBHEAD_HEIGHT,
        text="Solution",
        font_name=design.heading_font,
        font_size_pt=20,
        color_hex=design.accent_color,
        bold=True,
    )
    add_textbox(
        slide,
        left=BODY_LEFT,
        top=section2_top + _SUBHEAD_HEIGHT,
        width=_TEXT_WIDTH,
        height=Inches(1.5),
        text=content.proposed_solution_theme,
        font_name=design.body_font,
        font_size_pt=14,
        color_hex=design.text_color,
    )

    # ── Right-side image ────────────────────────────────────────────
    if blueprint.generated_image_path:
        img_region_left = Inches(7.6)
        img_region_top = Inches(1.6)
        img_max_w = Inches(5.2)
        img_max_h = Inches(4.8)
        try:
            w, h = fit_image_dims(
                blueprint.generated_image_path, img_max_w, img_max_h,
            )
            img_left = img_region_left + (img_max_w - w) // 2
            img_top = img_region_top + (img_max_h - h) // 2
            slide.shapes.add_picture(
                blueprint.generated_image_path,
                left=img_left, top=img_top, width=w, height=h,
            )
        except Exception:
            pass

    add_logo(slide, design)
    add_footer(slide, design)
