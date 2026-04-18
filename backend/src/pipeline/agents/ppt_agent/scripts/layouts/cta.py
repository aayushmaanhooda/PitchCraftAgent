"""cta layout — closing slide with centered next-steps + call-to-action.

Used for: the final slide of the deck (typically slide_type == 'closing'
on a main_slide entry). Reads title + key_bullets from that slide and
presents them as centered, emphasised next steps.

The title gets hero treatment (large, centred, primary_color). Bullets
are rendered as centred action items below. Footer is present as on
other slides.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.agents.ppt_agent.schema import (
    ContentSource,
    DensityVerdict,
    SalesPPTOutput,
    SlideBlueprint,
)
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    CONTENT_WIDTH,
    MARGIN_LEFT,
    add_footer,
    add_logo,
    add_textbox,
    hex_to_rgb,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design


def _get_content(
    content: SalesPPTOutput, blueprint: SlideBlueprint,
) -> tuple[str, list[str]]:
    """Extract title and action items."""
    src = blueprint.content_source

    if src == ContentSource.MAIN_SLIDE:
        slide = content.slides[blueprint.content_index]
        bullets = slide.key_bullets
        if (
            blueprint.density_verdict == DensityVerdict.SIMPLIFY
            and blueprint.simplified_bullets
        ):
            bullets = blueprint.simplified_bullets
        return slide.title, bullets

    # Fallback.
    return "Next Steps", []


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render a CTA slide — centred title + numbered action items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    title, bullets = _get_content(content, blueprint)

    # Hero title — centred, large.
    title_top = Inches(1.8)
    add_textbox(
        slide,
        left=MARGIN_LEFT,
        top=title_top,
        width=CONTENT_WIDTH,
        height=Inches(1.2),
        text=title,
        font_name=design.heading_font,
        font_size_pt=36,
        color_hex=design.primary_color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Numbered action items, left-aligned with indent.
    if bullets:
        items_top = title_top + Inches(1.6)
        items_left = MARGIN_LEFT + Inches(1.0)
        items_width = CONTENT_WIDTH - Inches(2.0)
        shape = slide.shapes.add_textbox(
            items_left, items_top, items_width, Inches(4.0),
        )
        tf = shape.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i + 1}.  {bullet}"
            p.font.name = design.body_font
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(design.text_color)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)

    add_logo(slide, design)
    add_footer(slide, design)
