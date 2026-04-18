"""icon_grid layout — title + bullet list where each bullet has a Font Awesome icon.

Used for: why-us slides, capability lists, anywhere a bullet list
benefits from visual anchors but doesn't warrant a generated image.

Icons come from blueprint.icons (list of FA names set by the visual
planner). Rasterised from pre-committed SVGs via icon_helpers at the
design's accent_color. Falls back to plain bullets if an icon is missing.
"""

from __future__ import annotations

import io

from pptx.util import Inches, Pt

from pipeline.agents.ppt_agent.schema import (
    ContentSource,
    DensityVerdict,
    SalesPPTOutput,
    SlideBlueprint,
)
from pipeline.agents.ppt_agent.scripts.render.icon_helpers import get_icon_png
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    BODY_LEFT,
    BODY_TOP,
    BODY_WIDTH,
    add_bullets,
    add_footer,
    add_logo,
    add_title,
    hex_to_rgb,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design

# Icon sizing and row spacing.
_ICON_SIZE = Inches(0.4)
_ICON_GAP = Inches(0.15)
_ROW_HEIGHT = Inches(0.75)


def _get_content(
    content: SalesPPTOutput, blueprint: SlideBlueprint,
) -> tuple[str, list[str]]:
    """Extract title and bullet list based on content_source."""
    src = blueprint.content_source

    if src == ContentSource.DECK_WHY_AZIRO:
        return "Why Aziro", content.why_aziro

    if src == ContentSource.DECK_DIFFERENTIATORS:
        return "Differentiators", content.differentiators

    if src == ContentSource.MAIN_SLIDE:
        slide = content.slides[blueprint.content_index]
        bullets = slide.key_bullets
        if (
            blueprint.density_verdict == DensityVerdict.SIMPLIFY
            and blueprint.simplified_bullets
        ):
            bullets = blueprint.simplified_bullets
        return slide.title, bullets

    # Fallback.
    return src.value.replace("_", " ").title(), []


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render an icon_grid slide — title + icon-prefixed bullet rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    title, bullets = _get_content(content, blueprint)
    add_title(slide, title, design)

    icons: list[str] = blueprint.icons or []

    # If no icons provided at all, fall back to plain bullets.
    if not icons:
        add_bullets(slide, bullets, design)
        add_logo(slide, design)
        add_footer(slide, design)
        return

    # Render each bullet as an icon + text row.
    text_left = BODY_LEFT + _ICON_SIZE + _ICON_GAP
    text_width = BODY_WIDTH - _ICON_SIZE - _ICON_GAP

    for i, bullet_text in enumerate(bullets):
        row_top = BODY_TOP + _ROW_HEIGHT * i

        # Try to render the icon.
        icon_name = icons[i] if i < len(icons) else None
        if icon_name:
            png_bytes = get_icon_png(icon_name, design.accent_color)
            if png_bytes:
                slide.shapes.add_picture(
                    io.BytesIO(png_bytes),
                    left=BODY_LEFT,
                    top=row_top,
                    width=_ICON_SIZE,
                    height=_ICON_SIZE,
                )

        # Bullet text next to the icon.
        shape = slide.shapes.add_textbox(
            text_left, row_top, text_width, _ROW_HEIGHT,
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet_text
        p.font.name = design.body_font
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(design.text_color)

    add_logo(slide, design)
    add_footer(slide, design)
