"""two_column layout — title + two equal columns side by side.

Primary use: use_case slides. Reads UseCase via
content.use_cases[blueprint.content_index]:
  Left column:  customer_problem + why_relevant_here.
  Right column: aziro_solution + business_value.

Also reusable for any main_slide if the planner picks this layout
(splits key_bullets evenly across two columns).

Title sized with heading_font + primary_color. Columns use body_font +
text_color. Column headers in accent_color for visual separation.
"""

from __future__ import annotations

from pptx.util import Inches, Pt

from pipeline.agents.ppt_agent.schema import (
    ContentSource,
    DensityVerdict,
    SalesPPTOutput,
    SlideBlueprint,
)
from pipeline.agents.ppt_agent.scripts.render.layout_helpers import (
    BODY_LEFT,
    BODY_TOP,
    BODY_WIDTH,
    add_footer,
    add_logo,
    add_textbox,
    add_title,
    hex_to_rgb,
    set_slide_bg,
)
from pipeline.agents.ppt_agent.scripts.theme.base_design import Design

_COL_GAP = Inches(0.6)
_COL_WIDTH = (BODY_WIDTH - _COL_GAP) // 2
_COL_HEADER_HEIGHT = Inches(0.55)
_COL_BODY_FONT_PT = 16
_SUBLABEL_FONT_PT = 14


def _add_column(
    slide,
    left,
    top,
    width,
    header: str,
    body_primary: str,
    sublabel: str | None,
    body_secondary: str | None,
    design: Design,
) -> None:
    """Render one column: header + primary prose + optional sub-labelled block."""
    # Coloured header band (accent colour, bold).
    add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=_COL_HEADER_HEIGHT,
        text=header,
        font_name=design.heading_font,
        font_size_pt=20,
        color_hex=design.accent_color,
        bold=True,
    )

    # Primary prose block (grows to absorb available height).
    body_top = top + _COL_HEADER_HEIGHT + Inches(0.15)
    primary_height = Inches(2.4)
    add_textbox(
        slide,
        left=left,
        top=body_top,
        width=width,
        height=primary_height,
        text=body_primary,
        font_name=design.body_font,
        font_size_pt=_COL_BODY_FONT_PT,
        color_hex=design.text_color,
    )

    if not (sublabel and body_secondary):
        return

    # Sub-label (bold, smaller, primary colour) + its prose below.
    sub_top = body_top + primary_height + Inches(0.15)
    add_textbox(
        slide,
        left=left,
        top=sub_top,
        width=width,
        height=Inches(0.35),
        text=sublabel,
        font_name=design.heading_font,
        font_size_pt=_SUBLABEL_FONT_PT,
        color_hex=design.primary_color,
        bold=True,
    )
    add_textbox(
        slide,
        left=left,
        top=sub_top + Inches(0.4),
        width=width,
        height=Inches(1.5),
        text=body_secondary,
        font_name=design.body_font,
        font_size_pt=_COL_BODY_FONT_PT,
        color_hex=design.text_color,
    )


def render(
    prs, content: SalesPPTOutput, blueprint: SlideBlueprint, design: Design,
) -> None:
    """Render a two_column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, design.background_color)

    src = blueprint.content_source
    left_col = BODY_LEFT
    right_col = BODY_LEFT + _COL_WIDTH + _COL_GAP

    if src == ContentSource.USE_CASE:
        uc = content.use_cases[blueprint.content_index]
        add_title(slide, uc.title, design)

        _add_column(
            slide,
            left=left_col,
            top=BODY_TOP,
            width=_COL_WIDTH,
            header="Challenge",
            body_primary=uc.customer_problem,
            sublabel="Why relevant",
            body_secondary=uc.why_relevant_here,
            design=design,
        )
        _add_column(
            slide,
            left=right_col,
            top=BODY_TOP,
            width=_COL_WIDTH,
            header="Aziro Solution",
            body_primary=uc.aziro_solution,
            sublabel="Business value",
            body_secondary=uc.business_value,
            design=design,
        )

    elif src == ContentSource.MAIN_SLIDE:
        ms = content.slides[blueprint.content_index]
        add_title(slide, ms.title, design)

        bullets = ms.key_bullets
        if (
            blueprint.density_verdict == DensityVerdict.SIMPLIFY
            and blueprint.simplified_bullets
        ):
            bullets = blueprint.simplified_bullets

        mid = len(bullets) // 2 or 1
        left_text = "\n\n".join(bullets[:mid])
        right_text = "\n\n".join(bullets[mid:])

        _add_column(
            slide, left_col, BODY_TOP, _COL_WIDTH,
            header="",
            body_primary=left_text,
            sublabel=None,
            body_secondary=None,
            design=design,
        )
        _add_column(
            slide, right_col, BODY_TOP, _COL_WIDTH,
            header="",
            body_primary=right_text,
            sublabel=None,
            body_secondary=None,
            design=design,
        )
    else:
        add_title(slide, src.value.replace("_", " ").title(), design)

    add_logo(slide, design)
    add_footer(slide, design)
