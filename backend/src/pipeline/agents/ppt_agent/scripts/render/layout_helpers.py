"""Shared positioning primitives for layout modules.

Keeps EMU arithmetic and python-pptx ceremony out of individual
layout files so they can focus on slide-specific arrangement.
"""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.agents.ppt_agent.scripts.theme.base_design import Design

# ---------------------------------------------------------------------------
# Slide dimensions — 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Standard margins and regions
# ---------------------------------------------------------------------------
MARGIN_LEFT = Inches(0.6)
MARGIN_RIGHT = Inches(0.6)
MARGIN_TOP = Inches(0.4)
MARGIN_BOTTOM = Inches(0.5)

CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

# Logo region reserved on the right of the title row.
LOGO_REGION_WIDTH = Inches(2.2)

# Title region — leaves room on the right for the logo so long titles
# don't collide with it.
TITLE_LEFT = MARGIN_LEFT
TITLE_TOP = MARGIN_TOP
TITLE_WIDTH = CONTENT_WIDTH - LOGO_REGION_WIDTH
TITLE_HEIGHT = Inches(1.0)

# Body region (below title)
BODY_TOP = TITLE_TOP + TITLE_HEIGHT + Inches(0.2)
BODY_LEFT = MARGIN_LEFT
BODY_WIDTH = CONTENT_WIDTH
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - MARGIN_BOTTOM

# Footer region
FOOTER_HEIGHT = Inches(0.45)
FOOTER_TOP = SLIDE_HEIGHT - FOOTER_HEIGHT

# Logo sizing
LOGO_HEIGHT = Inches(0.7)
LOGO_MARGIN = Inches(0.3)


# ---------------------------------------------------------------------------
# Colour helper
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a 6-char hex string (e.g. ``"1A3C6E"``) to an ``RGBColor``."""
    return RGBColor.from_string(hex_str)


# ---------------------------------------------------------------------------
# Slide background
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color_hex: str) -> None:
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


# ---------------------------------------------------------------------------
# Text primitives
# ---------------------------------------------------------------------------

def fit_heading_box(
    text: str,
    *,
    width_in: float,
    base_font_pt: int,
    min_font_pt: int = 22,
    max_lines: int = 3,
    line_spacing: float = 1.25,
    char_width_ratio: float = 0.62,
) -> tuple[int, float]:
    """Pick a font size + box height that keeps ``text`` from overflowing.

    Heuristic (no real text-measurement API in python-pptx):
      - Estimate chars-per-line at a given font size for a bold heading:
        avg char width ≈ ``char_width_ratio × font_pt`` in points.
        0.62 is tuned pessimistic — Georgia / Calibri Bold are both
        covered with room to spare so we err on the side of shrinking.
      - Subtract 0.2" for the text frame's internal padding (python-pptx
        default inset is 0.1" per side).
      - Account for word-wrap waste: real wrap boundaries land on spaces,
        which typically wastes ~15% of each line vs pure char-based math.
      - If the text wraps to more than ``max_lines`` at ``base_font_pt``,
        step the font down by 2pt and retry, until ``min_font_pt``.
      - Return the final font size AND the box height (inches) sized to
        the actual wrapped-line count so the next element can sit just
        below the box instead of at a hardcoded offset.

    Not pixel-perfect, but reliably prevents the overlap seen on long
    cover titles.
    """
    import math

    length = len(text)
    effective_width_pt = max(1.0, (width_in - 0.2) * 72)
    font_pt = base_font_pt

    while font_pt >= min_font_pt:
        chars_per_line = max(
            1, int(effective_width_pt / (char_width_ratio * font_pt) * 0.85),
        )
        lines = max(1, math.ceil(length / chars_per_line))
        if lines <= max_lines:
            break
        font_pt -= 2
    else:
        font_pt = min_font_pt
        chars_per_line = max(
            1, int(effective_width_pt / (char_width_ratio * font_pt) * 0.85),
        )
        lines = max(1, math.ceil(length / chars_per_line))

    height_in = (lines * font_pt * line_spacing) / 72 + 0.25
    return font_pt, height_in


def add_textbox(slide, left, top, width, height, text: str,
                font_name: str, font_size_pt: int, color_hex: str,
                bold: bool = False, alignment=PP_ALIGN.LEFT,
                word_wrap: bool = True):
    """Add a text box with a single styled paragraph.

    Returns ``(shape, text_frame)`` so callers can append more paragraphs
    if needed.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size_pt)
    p.font.color.rgb = hex_to_rgb(color_hex)
    p.font.bold = bold
    p.alignment = alignment
    return shape, tf


def add_title(slide, text: str, design: Design, *,
              left=None, top=None, width=None, height=None,
              font_size_pt: int = 32,
              alignment=PP_ALIGN.LEFT) -> None:
    """Add the slide title using the design's heading font + primary colour."""
    add_textbox(
        slide,
        left=left if left is not None else TITLE_LEFT,
        top=top if top is not None else TITLE_TOP,
        width=width if width is not None else TITLE_WIDTH,
        height=height if height is not None else TITLE_HEIGHT,
        text=text,
        font_name=design.heading_font,
        font_size_pt=font_size_pt,
        color_hex=design.primary_color,
        bold=True,
        alignment=alignment,
    )


# ---------------------------------------------------------------------------
# Numbered list helper
# ---------------------------------------------------------------------------

def _add_numbered_paragraph(tf, index: int, text: str, design: Design,
                            font_size_pt: int, space_after_pt: int,
                            first: bool = False) -> None:
    """Append a paragraph with a bold accent-coloured number + body text.

    Uses two runs so the number (e.g. ``"1."``) stands out visually while
    the text uses the design's normal body colour. Numbered prefixes are
    more reliable than OoXML bullets in python-pptx blank textboxes.
    """
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    # Clear any default text state.
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after_pt)

    run_num = p.add_run()
    run_num.text = f"{index}.  "
    run_num.font.name = design.heading_font
    run_num.font.size = Pt(font_size_pt)
    run_num.font.bold = True
    run_num.font.color.rgb = hex_to_rgb(design.accent_color)

    run_txt = p.add_run()
    run_txt.text = text
    run_txt.font.name = design.body_font
    run_txt.font.size = Pt(font_size_pt)
    run_txt.font.color.rgb = hex_to_rgb(design.text_color)


def add_bullets(slide, bullets: list[str], design: Design, *,
                left=None, top=None, width=None, height=None,
                font_size_pt: int = 18,
                bullet_char: str | None = None,  # kept for API compat; unused
                space_after_pt: int = 10) -> None:
    """Render a list as numbered items (1., 2., 3., ...).

    The numeric prefix is bold + accent-coloured, the body text is
    design.text_color. We use explicit runs rather than OoXML bullet
    chars because blank python-pptx textboxes don't always honour
    buChar — numbers always render.
    """
    shape = slide.shapes.add_textbox(
        left if left is not None else BODY_LEFT,
        top if top is not None else BODY_TOP,
        width if width is not None else BODY_WIDTH,
        height if height is not None else BODY_HEIGHT,
    )
    tf = shape.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        _add_numbered_paragraph(
            tf,
            index=i + 1,
            text=bullet,
            design=design,
            font_size_pt=font_size_pt,
            space_after_pt=space_after_pt,
            first=(i == 0),
        )


# ---------------------------------------------------------------------------
# Body text
# ---------------------------------------------------------------------------

def add_body(slide, text: str, design: Design, *,
             left=None, top=None, width=None, height=None,
             font_size_pt: int = 16) -> None:
    """Add a body paragraph styled with the design's body font + text colour."""
    add_textbox(
        slide,
        left=left if left is not None else BODY_LEFT,
        top=top if top is not None else BODY_TOP,
        width=width if width is not None else BODY_WIDTH,
        height=height if height is not None else BODY_HEIGHT,
        text=text,
        font_name=design.body_font,
        font_size_pt=font_size_pt,
        color_hex=design.text_color,
    )


# ---------------------------------------------------------------------------
# Slide chrome (footer + logo)
# ---------------------------------------------------------------------------

def add_footer(slide, design: Design) -> None:
    """Add footer text at the bottom of the slide."""
    add_textbox(
        slide,
        left=MARGIN_LEFT,
        top=FOOTER_TOP,
        width=CONTENT_WIDTH,
        height=FOOTER_HEIGHT,
        text=design.footer_text,
        font_name=design.body_font,
        font_size_pt=12,
        color_hex=design.secondary_color,
        alignment=PP_ALIGN.LEFT,
    )


def add_logo(slide, design: Design) -> None:
    """Add the brand logo to the top-right corner of the slide.

    Aspect ratio is preserved — only height is fixed, width auto-scales.
    The logo is then repositioned to sit flush against the right margin.
    """
    if design.logo_path is None or not Path(design.logo_path).exists():
        return
    pic = slide.shapes.add_picture(
        str(design.logo_path),
        left=0,
        top=LOGO_MARGIN,
        height=LOGO_HEIGHT,
    )
    # Right-align: reposition now that we know the auto-scaled width.
    pic.left = SLIDE_WIDTH - pic.width - LOGO_MARGIN
